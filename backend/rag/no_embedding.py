import PyPDF2
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from docx import Document
from pptx import Presentation
import json
import threading
from concurrent.futures import ThreadPoolExecutor  # Added for parallel PDF processing

# Cache globals and lock
CACHE_FILE = os.path.join(os.path.dirname(__file__), "extracted_cache.json")
extracted_cache = {}
cache_lock = threading.Lock()
last_file_path = None

# Load cache from file if it exists
if os.path.exists(CACHE_FILE):
    try:
        with open(CACHE_FILE, "r", encoding="utf-8") as f:
            extracted_cache = json.load(f)
    except Exception as e:
        print("Failed to load cache:", e)


def save_cache():
    with cache_lock:
        try:
            with open(CACHE_FILE, "w", encoding="utf-8") as f:
                json.dump(extracted_cache, f)
        except Exception as e:
            print("Failed to save cache:", e)


def get_cached_extracted_text(file_path):
    global last_file_path
    try:
        file_mod_time = os.path.getmtime(file_path)
    except Exception as e:
        print("Failed to get modification time:", e)
        file_mod_time = None
    # Only use cache if file is same as previous
    if last_file_path == file_path:
        cache_entry = extracted_cache.get(file_path)
        if cache_entry and cache_entry.get("mod_time") == file_mod_time:
            return cache_entry.get("text", "")
    # If new file or no valid cache, extract fresh and update last_file_path.
    text = extract_text_from_file(file_path)
    if last_file_path in extracted_cache:
        del extracted_cache[last_file_path]
    extracted_cache[file_path] = {"mod_time": file_mod_time, "text": text}
    last_file_path = file_path
    save_cache()
    return text


def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        with ThreadPoolExecutor() as executor:
            page_texts = list(
                executor.map(lambda page: page.extract_text() or "", reader.pages)
            )
        text = "\n".join(page_texts)
    return text


def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading docx file {file_path}: {e}")
    return text


def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading pptx file {file_path}: {e}")
    return text


def extract_text_from_plain(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    except Exception as e:
        print(f"Error reading plain text file {file_path}: {e}")
    return text


def extract_text_from_file(file_path):
    _, file_extension = os.path.splitext(file_path.lower())
    text = ""

    if file_extension == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif file_extension == ".docx":
        text = extract_text_from_docx(file_path)
    elif file_extension == ".pptx":
        text = extract_text_from_pptx(file_path)
    elif file_extension in [
        ".txt",
        ".md",
    ]:
        text = extract_text_from_plain(file_path)
    else:
        print(f"Unsupported file type: {file_extension}")

    return text


def split_text_to_chunks(text, chunk_size=1000, chunk_overlap=100):
    text_splitter = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""],
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        is_separator_regex=False,
    )
    chunks = text_splitter.split_text(text)
    return chunks


def keyword_search(query, documents, top_k=3):
    query_tokens = set(query.lower().split())
    scored_docs = []
    for doc in documents:
        doc_tokens = set(doc.lower().split())
        common_tokens = query_tokens.intersection(doc_tokens)
        score = len(common_tokens)
        if score > 0:
            scored_docs.append((score, doc))
    scored_docs.sort(reverse=True, key=lambda x: x[0])
    return [doc for _, doc in scored_docs[:top_k]]


def search_file_for_keywords(
    file_path, query, top_k=3, chunk_size=1000, chunk_overlap=100
):
    # Use cache to retrieve text if available
    text = get_cached_extracted_text(file_path)
    if not text:
        return []

    coding_extensions = [
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".c",
        ".cpp",
        ".java",
        ".html",
        ".css",
    ]
    if os.path.splitext(file_path.lower())[1] in coding_extensions:
        return [text]
    else:
        chunks = split_text_to_chunks(text, chunk_size, chunk_overlap)
        return keyword_search(query, chunks, top_k)
